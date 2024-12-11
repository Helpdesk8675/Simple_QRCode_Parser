import os
import io
import tkinter as tk
from tkinter import filedialog, messagebox
import pandas as pd
from pyzbar.pyzbar import decode
from PIL import Image
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook

def extract_images_from_docx(file_path):
    images = []
    try:
        doc = Document(file_path)
        rels = doc.part.rels
        for rel in rels.values():
            if "image" in rel.target_ref:
                image_blob = rel.target_part.blob
                images.append(image_blob)
    except Exception as e:
        print(f"Error processing DOCX: {e}")
    return images

def extract_images_from_pptx(file_path):
    images = []
    try:
        ppt = Presentation(file_path)
        for slide in ppt.slides:
            for shape in slide.shapes:
                if shape.shape_type == 13:  # Picture
                    images.append(shape.image.blob)
    except Exception as e:
        print(f"Error processing PPTX: {e}")
    return images

def extract_images_from_xlsx(file_path):
    images = []
    try:
        wb = load_workbook(file_path)
        for sheet in wb.worksheets:
            for image in sheet._images:
                images.append(image._data)
    except Exception as e:
        print(f"Error processing XLSX: {e}")
    return images

def parse_qrcodes_in_folder(source_folder, output_folder):
    results = []

    for root, _, files in os.walk(source_folder):
        for file in files:
            file_path = os.path.join(root, file)
            try:
                if file.lower().endswith((".png", ".jpg", ".jpeg", ".bmp", ".gif")):
                    image = Image.open(file_path)
                    qr_codes = decode(image)
                    for qr_code in qr_codes:
                        url = qr_code.data.decode('utf-8')
                        results.append({"File Name": file, "URL": url})
                elif file.lower().endswith(".docx"):
                    images = extract_images_from_docx(file_path)
                    for img_data in images:
                        image = Image.open(io.BytesIO(img_data))
                        qr_codes = decode(image)
                        for qr_code in qr_codes:
                            url = qr_code.data.decode('utf-8')
                            results.append({"File Name": file, "URL": url})
                elif file.lower().endswith(".pptx"):
                    images = extract_images_from_pptx(file_path)
                    for img_data in images:
                        image = Image.open(io.BytesIO(img_data))
                        qr_codes = decode(image)
                        for qr_code in qr_codes:
                            url = qr_code.data.decode('utf-8')
                            results.append({"File Name": file, "URL": url})
                elif file.lower().endswith(".xlsx"):
                    images = extract_images_from_xlsx(file_path)
                    for img_data in images:
                        image = Image.open(io.BytesIO(img_data))
                        qr_codes = decode(image)
                        for qr_code in qr_codes:
                            url = qr_code.data.decode('utf-8')
                            results.append({"File Name": file, "URL": url})
            except Exception as e:
                print(f"Could not process file: {file_path}, Error: {e}")

    if results:
        output_file = os.path.join(output_folder, 'qrcode_results.csv')
        df = pd.DataFrame(results)
        df.to_csv(output_file, index=False)
        messagebox.showinfo("Success", f"QR Code data saved to {output_file}")
    else:
        messagebox.showinfo("No QR Codes", "No QR codes found in the specified folder.")

def select_source_folder():
    folder = filedialog.askdirectory(title="Select Source Folder")
    source_folder_var.set(folder)

def select_output_folder():
    folder = filedialog.askdirectory(title="Select Output Folder")
    output_folder_var.set(folder)

def start_processing():
    source_folder = source_folder_var.get()
    output_folder = output_folder_var.get()

    if not source_folder or not os.path.isdir(source_folder):
        messagebox.showerror("Error", "Please select a valid source folder.")
        return

    if not output_folder or not os.path.isdir(output_folder):
        messagebox.showerror("Error", "Please select a valid output folder.")
        return

    parse_qrcodes_in_folder(source_folder, output_folder)

# Create GUI
root = tk.Tk()
root.title("QR Code Extractor")

# Source folder selection
source_folder_var = tk.StringVar()
tk.Label(root, text="Source Folder:").grid(row=0, column=0, padx=10, pady=10, sticky="e")
tk.Entry(root, textvariable=source_folder_var, width=50).grid(row=0, column=1, padx=10, pady=10)
tk.Button(root, text="Browse", command=select_source_folder).grid(row=0, column=2, padx=10, pady=10)

# Output folder selection
output_folder_var = tk.StringVar()
tk.Label(root, text="Output Folder:").grid(row=1, column=0, padx=10, pady=10, sticky="e")
tk.Entry(root, textvariable=output_folder_var, width=50).grid(row=1, column=1, padx=10, pady=10)
tk.Button(root, text="Browse", command=select_output_folder).grid(row=1, column=2, padx=10, pady=10)

# Start button
tk.Button(root, text="Start", command=start_processing).grid(row=2, column=1, pady=20)

root.mainloop()
