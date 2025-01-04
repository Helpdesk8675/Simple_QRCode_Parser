import os
import io
import tkinter as tk
from tkinter import filedialog, messagebox, ttk
import pandas as pd
from pyzbar.pyzbar import decode
from PIL import Image
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook

class QRCodeParser:
    def __init__(self, root):
        self.root = root
        self.setup_gui()

    def extract_images_from_docx(self, file_path):
        images = []
        try:
            doc = Document(file_path)
            rels = doc.part.rels
            for rel in rels.values():
                if "image" in rel.target_ref:
                    image_blob = rel.target_part.blob
                    images.append(image_blob)
        except Exception as e:
            print(f"Error processing DOCX: {e}")
        return images

    def extract_images_from_pptx(self, file_path):
        images = []
        try:
            ppt = Presentation(file_path)
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if shape.shape_type == 13:  # Picture
                        images.append(shape.image.blob)
        except Exception as e:
            print(f"Error processing PPTX: {e}")
        return images

    def extract_images_from_xlsx(self, file_path):
        images = []
        try:
            wb = load_workbook(file_path)
            for sheet in wb.worksheets:
                for image in sheet._images:
                    images.append(image._data)
        except Exception as e:
            print(f"Error processing XLSX: {e}")
        return images

    def parse_qrcodes_in_folder(self, source_folder, output_folder):
        results = []
        
        # Get total number of files for progress calculation
        total_files = sum([len(files) for _, _, files in os.walk(source_folder)])
        processed_files = 0
        
        for root, _, files in os.walk(source_folder):
            for file in files:
                file_path = os.path.join(root, file)
                try:
                    # Update status label
                    self.status_label.config(text=f"Processing: {file}")
                    self.root.update_idletasks()

                    if file.lower().endswith((".png", ".jpg", ".jpeg", ".bmp", ".gif")):
                        image = Image.open(file_path)
                        qr_codes = decode(image)
                        for qr_code in qr_codes:
                            url = qr_code.data.decode('utf-8')
                            results.append({"File Name": file, "URL": url})
                    
                    elif file.lower().endswith(".docx"):
                        images = self.extract_images_from_docx(file_path)
                        for img_data in images:
                            image = Image.open(io.BytesIO(img_data))
                            qr_codes = decode(image)
                            for qr_code in qr_codes:
                                url = qr_code.data.decode('utf-8')
                                results.append({"File Name": file, "URL": url})
                    
                    elif file.lower().endswith(".pptx"):
                        images = self.extract_images_from_pptx(file_path)
                        for img_data in images:
                            image = Image.open(io.BytesIO(img_data))
                            qr_codes = decode(image)
                            for qr_code in qr_codes:
                                url = qr_code.data.decode('utf-8')
                                results.append({"File Name": file, "URL": url})
                    
                    elif file.lower().endswith(".xlsx"):
                        images = self.extract_images_from_xlsx(file_path)
                        for img_data in images:
                            image = Image.open(io.BytesIO(img_data))
                            qr_codes = decode(image)
                            for qr_code in qr_codes:
                                url = qr_code.data.decode('utf-8')
                                results.append({"File Name": file, "URL": url})
                
                except Exception as e:
                    print(f"Could not process file: {file_path}, Error: {e}")
                
                # Update progress
                processed_files += 1
                progress = (processed_files / total_files) * 100
                self.progress_var.set(progress)
                self.root.update_idletasks()

        # Update status for saving results
        self.status_label.config(text="Saving results...")
        self.root.update_idletasks()

        if results:
            output_file = os.path.join(output_folder, 'qrcode_results.csv')
            df = pd.DataFrame(results)
            df.to_csv(output_file, index=False)
            self.status_label.config(text="Complete!")
            messagebox.showinfo("Success", f"QR Code data saved to {output_file}")
        else:
            self.status_label.config(text="No QR codes found")
            messagebox.showinfo("No QR Codes", "No QR codes found in the specified folder.")

    def select_source_folder(self):
        folder = filedialog.askdirectory(title="Select Source Folder")
        self.source_folder_var.set(folder)

    def select_output_folder(self):
        folder = filedialog.askdirectory(title="Select Output Folder")
        self.output_folder_var.set(folder)

    def start_processing(self):
        source_folder = self.source_folder_var.get()
        output_folder = self.output_folder_var.get()

        if not source_folder or not os.path.isdir(source_folder):
            messagebox.showerror("Error", "Please select a valid source folder.")
            return

        if not output_folder or not os.path.isdir(output_folder):
            messagebox.showerror("Error", "Please select a valid output folder.")
            return

        # Disable the start button during processing
        self.start_button.config(state='disabled')
        
        # Create progress bar
        self.progress_var.set(0)
        self.progress_bar.grid(row=3, column=0, columnspan=3, padx=10, pady=10, sticky="ew")
        self.status_label.grid(row=4, column=0, columnspan=3, padx=10, pady=5)
        
        self.parse_qrcodes_in_folder(source_folder, output_folder)
        
        # Re-enable the start button
        self.start_button.config(state='normal')

    def setup_gui(self):
        self.root.title("QR Code Extractor by helpdesk8675")

        # Source folder selection
        self.source_folder_var = tk.StringVar()
        tk.Label(self.root, text="Source Folder:").grid(row=0, column=0, padx=10, pady=10, sticky="e")
        tk.Entry(self.root, textvariable=self.source_folder_var, width=50).grid(row=0, column=1, padx=10, pady=10)
        tk.Button(self.root, text="Browse", command=self.select_source_folder).grid(row=0, column=2, padx=10, pady=10)

        # Output folder selection
        self.output_folder_var = tk.StringVar()
        tk.Label(self.root, text="Output Folder:").grid(row=1, column=0, padx=10, pady=10, sticky="e")
        tk.Entry(self.root, textvariable=self.output_folder_var, width=50).grid(row=1, column=1, padx=10, pady=10)
        tk.Button(self.root, text="Browse", command=self.select_output_folder).grid(row=1, column=2, padx=10, pady=10)

        # Start button
        self.start_button = tk.Button(self.root, text="Start", command=self.start_processing)
        self.start_button.grid(row=2, column=1, pady=20)

        # Progress bar and status label (hidden initially)
        self.progress_var = tk.DoubleVar()
        self.progress_bar = ttk.Progressbar(self.root, variable=self.progress_var, maximum=100)
        self.status_label = tk.Label(self.root, text="")

def main():
    root = tk.Tk()
    app = QRCodeParser(root)
    root.mainloop()

if __name__ == "__main__":
    main()
